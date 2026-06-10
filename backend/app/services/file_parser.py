import io
import logging

logger = logging.getLogger("finswarm.file_parser")

def extract_pdf_text(file_bytes: bytes) -> str:
    """Extracts text from PDFs, preserving layout where possible and handling locks."""
    try:
        from pypdf import PdfReader
        
        stream = io.BytesIO(file_bytes)
        reader = PdfReader(stream)
        
        # Unlock standard read-only locks
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception as e:
                logger.warning(f"PDF encrypted: {e}")
                return "Error: PDF is encrypted and requires a password."

        full_text = []
        for idx, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text(extraction_mode="layout")
            except Exception:
                page_text = page.extract_text()
                
            if page_text:
                page_text = page_text.replace("\x00", " ") # Prevent JSON serialization errors
                full_text.append(f"--- PAGE {idx+1} ---\n{page_text.strip()}")
                
        if not full_text:
            return "Notice: PDF parsed successfully, but no text found (likely a scanned image)."
            
        return "\n\n".join(full_text)
        
    except ImportError:
        logger.error("pypdf missing.")
        return "Error: pypdf library is not installed (run `pip install pypdf`)."
    except Exception as e:
        logger.exception(f"PDF Parse Error: {e}")
        return f"Error extracting PDF: {str(e)}"

def extract_docx_text(file_bytes: bytes) -> str:
    """Extracts plain text from Microsoft Word documents."""
    try:
        import docx
        stream = io.BytesIO(file_bytes)
        doc = docx.Document(stream)
        
        full_text = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(full_text)
        
    except ImportError:
        logger.error("python-docx missing.")
        return "Error: python-docx library is not installed (run `pip install python-docx`)."
    except Exception as e:
        logger.exception(f"DOCX Parse Error: {e}")
        return f"Error extracting Word Document: {str(e)}"

def extract_excel_text(file_bytes: bytes) -> str:
    """Extracts text data from Excel spreadsheets, organizing it by sheet and row."""
    try:
        import openpyxl
        stream = io.BytesIO(file_bytes)
        # data_only=True extracts the calculated values rather than the raw formulas
        wb = openpyxl.load_workbook(filename=stream, data_only=True)
        
        full_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            full_text.append(f"\n--- SHEET: {sheet_name} ---")
            
            for row in ws.iter_rows(values_only=True):
                # Filter out empty cells and join the row with tabs
                row_vals = [str(cell) for cell in row if cell is not None]
                if row_vals:
                    full_text.append(" | ".join(row_vals))
                    
        return "\n".join(full_text)
        
    except ImportError:
        logger.error("openpyxl missing.")
        return "Error: openpyxl library is not installed (run `pip install openpyxl`)."
    except Exception as e:
        logger.exception(f"Excel Parse Error: {e}")
        return f"Error extracting Excel File: {str(e)}"

def extract_ppt_text(file_bytes: bytes) -> str:
    """Extracts text shapes from PowerPoint presentations, organizing by slide."""
    try:
        from pptx import Presentation
        stream = io.BytesIO(file_bytes)
        prs = Presentation(stream)
        
        full_text = []
        for i, slide in enumerate(prs.slides):
            full_text.append(f"\n--- SLIDE {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text.strip())
                    
        return "\n".join(full_text)
        
    except ImportError:
        logger.error("python-pptx missing.")
        return "Error: python-pptx library is not installed (run `pip install python-pptx`)."
    except Exception as e:
        logger.exception(f"PPTX Parse Error: {e}")
        return f"Error extracting PowerPoint: {str(e)}"